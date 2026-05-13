import logging
from pathlib import Path
from typing import List

from pptx import Presentation

from ingestion.types import PageUnit

logger = logging.getLogger(__name__)


def load_pptx(file_path: Path) -> List[PageUnit]:
    logger.info("Loading PPTX: %s", file_path)
    presentation = Presentation(str(file_path))
    units: List[PageUnit] = []

    for i, slide in enumerate(presentation.slides, start=1):
        text_blocks: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = (shape.text or "").strip()
                if text:
                    text_blocks.append(text)

        if not text_blocks:
            continue

        units.append(
            PageUnit(
                id=f"{file_path.name}:slide:{i}",
                content="\n".join(text_blocks),
                unit_type="slide",
                index=i,
                source=file_path.name,
            )
        )

    if not units:
        raise ValueError(f"No text extracted from PPTX: {file_path}")

    return units
